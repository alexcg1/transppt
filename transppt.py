#!/usr/bin/python3

from os import environ

# You'll need to use your Google API key for the translate service
API_KEY = "Google_API.json"
environ["GOOGLE_APPLICATION_CREDENTIALS"] = API_KEY

import argparse
from pptx import Presentation
from google.cloud import translate

try:
	from progress.bar import ShadyBar
	progress_check = True
except:
	progress_check = False
	pass

def zh_to_en(text):
	translation = translate_client.translate(text, target_language=target)
	return translation['translatedText']

translate_client = translate.Client()
target = "en"

# Set up command line arguments
parser = argparse.ArgumentParser()
parser.add_argument("input_file")
args = parser.parse_args()
parser.parse_args()

# Open presentation
print("Opening", args.input_file)
input_ppt = Presentation(args.input_file)

output_ppt = input_ppt # Create output presentation

i = 0 # Set up iterator to track slides between input and output files
element_count = 0 # Set up iterator to track number of elements translated

# Set up progress bar
if progress_check:
	slide_count = 0
	for slide in output_ppt.slides:
		slide_count += 1
	bar = ShadyBar('Translating', max=slide_count)

# Translate slides
for slide in output_ppt.slides:
	if progress_check == False:
		print("\tTranslating slide",i)
	j = 0 # Set up iterator j to go over each item in slide.shapes array
	for shape in slide.shapes:
		try:
			shape.text = zh_to_en(input_ppt.slides[i].shapes[j].text) # This SHOULD work, right?
			element_count += 1
			j += 1
		except:
			pass
	i += 1
	bar.next()

# Write presentation to new file
output_file = target+"_"+args.input_file
print("\nSaving as",output_file)
output_ppt.save(output_file)
bar.finish()

# print(element_count, "elements of presentation translated")